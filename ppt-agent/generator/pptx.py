from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import io
import base64
from PIL import Image
import requests
from lxml import etree

# 颜色工具
def hex_to_rgb(hex_color: str) -> tuple:
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def rgb_to_hex(r: int, g: int, b: int) -> str:
    return f"{r:02X}{g:02X}{b:02X}"


def blend_colors(color1: tuple, color2: tuple, ratio: float) -> tuple:
    """混合两个颜色"""
    r = int(color1[0] * (1 - ratio) + color2[0] * ratio)
    g = int(color1[1] * (1 - ratio) + color2[1] * ratio)
    b = int(color1[2] * (1 - ratio) + color2[2] * ratio)
    return (r, g, b)


# 主构建函数
def build_pptx(content: dict, style: dict, images: dict = None) -> str:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    primary = hex_to_rgb(style["primary"])
    secondary = hex_to_rgb(style["secondary"])
    accent = hex_to_rgb(style["accent"])
    highlight = hex_to_rgb(style.get("highlight", "E53E3E"))

    layouts = style.get("layouts", ["center"])

    for i, slide_data in enumerate(content.get("slides", [])):
        slide_type = slide_data.get("type", "content")
        layout = layouts[i % len(layouts)]

        if slide_type == "cover":
            add_cover_slide(prs, slide_data, primary, secondary, accent, style)
        elif slide_type == "content":
            if layout == "split":
                add_split_slide(prs, slide_data, primary, secondary, accent, style, images)
            elif layout == "grid":
                add_grid_slide(prs, slide_data, primary, secondary, accent, style, images)
            elif layout == "cards":
                add_cards_slide(prs, slide_data, primary, secondary, accent, style, images)
            elif layout == "two-col":
                add_two_col_slide(prs, slide_data, primary, secondary, accent, style, images)
            else:
                add_content_slide(prs, slide_data, primary, secondary, accent, style)
        elif slide_type == "ending":
            add_ending_slide(prs, slide_data, primary, secondary, accent, style)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return base64.b64encode(output.read()).decode("utf-8")


def set_shape_gradient(shape, color1: tuple, color2: tuple, angle: int = 0):
    """设置形状渐变填充"""
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = angle
    fill.gradient_stops[0].color.rgb = RGBColor(*color1)
    fill.gradient_stops[1].color.rgb = RGBColor(*color2)


def add_rectangle(slide, x, y, w, h, color: tuple, alpha: float = 1.0):
    """添加矩形"""
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*color)
    shape.line.fill.background()
    return shape


def add_decorative_circle(slide, x, y, size, color: tuple, alpha: float = 0.1):
    """添加装饰圆形"""
    shape = slide.shapes.add_shape(9, x, y, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*color)
    shape.line.fill.background()
    return shape


def add_image_to_slide(slide, url: str, x, y, w, h):
    """添加网络图片到幻灯片"""
    try:
        resp = requests.get(url, timeout=10)
        if resp.status_code == 200:
            img_data = io.BytesIO(resp.content)
            slide.shapes.add_picture(img_data, x, y, w, h)
            return True
    except Exception:
        pass
    return False


def add_text(slide, text, x, y, w, h, font_size, color: tuple, bold: bool = False,
             font_name: str = "微软雅黑", align: int = PP_ALIGN.LEFT, line_spacing: float = 1.5):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(*color)
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    p.line_spacing = line_spacing
    return txBox


# 封面页 - 深色渐变背景 + 大标题
def add_cover_slide(prs, data, primary, secondary, accent, style):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 渐变背景
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_gradient(bg, primary, secondary, 135)
    bg.line.fill.background()

    # 装饰圆形
    add_decorative_circle(slide, Inches(9), Inches(-1), Inches(5), secondary, 0.15)
    add_decorative_circle(slide, Inches(-1), Inches(5), Inches(4), accent, 0.08)

    # 左侧色块装饰
    accent_bar = add_rectangle(slide, 0, Inches(2.8), Inches(0.15), Inches(1.5), accent)

    # 主标题
    add_text(slide, data.get("title", ""),
             Inches(0.8), Inches(2.5), Inches(10), Inches(1.2),
             48, accent, True, style["font_title"], PP_ALIGN.LEFT)

    # 副标题
    add_text(slide, data.get("subtitle", ""),
             Inches(0.8), Inches(3.8), Inches(9), Inches(0.8),
             24, (255, 255, 255), False, style["font_body"], PP_ALIGN.LEFT)

    # 底部装饰线
    line = add_rectangle(slide, Inches(0.8), Inches(4.8), Inches(3), Inches(0.03), accent)


# 内容页 - 标准布局
def add_content_slide(prs, data, primary, secondary, accent, style):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 顶部色带
    add_rectangle(slide, 0, 0, prs.slide_width, Inches(0.08), accent)
    header = add_rectangle(slide, 0, 0, prs.slide_width, Inches(1.4), primary)

    # 页面标题
    add_text(slide, data.get("title", ""),
             Inches(0.5), Inches(0.35), Inches(12), Inches(0.9),
             32, (255, 255, 255), True, style["font_title"], PP_ALIGN.LEFT)

    # 左侧强调色块
    add_rectangle(slide, 0, Inches(1.4), Inches(0.08), Inches(5.5), accent)

    # 内容区域
    bullets = data.get("bullets", [])
    content = "\n".join([f"• {b}" for b in bullets])
    add_text(slide, content,
             Inches(0.6), Inches(1.8), Inches(12), Inches(5),
             20, (51, 51, 51), False, style["font_body"], PP_ALIGN.LEFT, 1.8)


# 分栏布局 - 左侧文字右侧图
def add_split_slide(prs, data, primary, secondary, accent, style, images):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 左侧深色区域
    add_rectangle(slide, 0, 0, Inches(5.5), prs.slide_height, primary)
    add_decorative_circle(slide, Inches(-2), Inches(4), Inches(6), secondary, 0.2)

    # 左侧标题和内容
    add_text(slide, data.get("title", ""),
             Inches(0.4), Inches(1.5), Inches(4.8), Inches(1),
             28, accent, True, style["font_title"], PP_ALIGN.LEFT)

    bullets = data.get("bullets", [])
    content = "\n".join([f"• {b}" for b in bullets])
    add_text(slide, content,
             Inches(0.4), Inches(2.6), Inches(4.8), Inches(4),
             18, (230, 230, 230), False, style["font_body"], PP_ALIGN.LEFT, 1.6)

    # 右侧图片区
    img_url = images.get(data.get("title", "")) if images else None
    if img_url:
        add_image_to_slide(slide, img_url, Inches(5.8), Inches(1), Inches(7), Inches(5.5))
    else:
        # 占位背景
        add_rectangle(slide, Inches(5.8), Inches(1), Inches(7), Inches(5.5), (240, 240, 245))
        add_text(slide, "图片区",
                 Inches(5.8), Inches(3), Inches(7), Inches(1),
                 24, (180, 180, 180), False, style["font_title"], PP_ALIGN.CENTER)


# 网格布局 - 2x2卡片
def add_grid_slide(prs, data, primary, secondary, accent, style, images):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 顶部色带
    add_rectangle(slide, 0, 0, prs.slide_width, Inches(0.08), accent)
    add_rectangle(slide, 0, 0, prs.slide_width, Inches(1.3), primary)

    add_text(slide, data.get("title", ""),
             Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             30, (255, 255, 255), True, style["font_title"], PP_ALIGN.LEFT)

    # 2x2网格卡片
    bullets = data.get("bullets", [])
    card_w = Inches(5.8)
    card_h = Inches(2.5)
    start_x = Inches(0.5)
    start_y = Inches(1.6)
    gap = Inches(0.4)

    for idx, bullet in enumerate(bullets[:4]):
        col = idx % 2
        row = idx // 2
        x = start_x + col * (card_w + gap)
        y = start_y + row * (card_h + gap)

        # 卡片背景
        card = add_rectangle(slide, x, y, card_w, card_h, (248, 248, 250))
        # 顶部强调线
        add_rectangle(slide, x, y, card_w, Inches(0.06), accent)

        # 序号
        add_text(slide, f"0{idx + 1}",
                 x + Inches(0.15), y + Inches(0.15), Inches(0.6), Inches(0.5),
                 20, accent, True, style["font_title"], PP_ALIGN.LEFT)

        # 内容
        add_text(slide, bullet,
                 x + Inches(0.15), y + Inches(0.6), card_w - Inches(0.3), card_h - Inches(0.8),
                 16, (60, 60, 60), False, style["font_body"], PP_ALIGN.LEFT, 1.4)


# 卡片布局 - 横向排列
def add_cards_slide(prs, data, primary, secondary, accent, style, images):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 顶部色带
    add_rectangle(slide, 0, 0, prs.slide_width, Inches(0.08), accent)
    add_rectangle(slide, 0, 0, prs.slide_width, Inches(1.3), primary)

    add_text(slide, data.get("title", ""),
             Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             30, (255, 255, 255), True, style["font_title"], PP_ALIGN.LEFT)

    # 横向3卡片
    bullets = data.get("bullets", [])
    card_w = Inches(3.8)
    card_h = Inches(4.8)
    start_x = Inches(0.5)
    start_y = Inches(1.6)
    gap = Inches(0.5)

    for idx, bullet in enumerate(bullets[:3]):
        x = start_x + idx * (card_w + gap)
        y = start_y

        # 卡片背景 - 白色
        add_rectangle(slide, x, y, card_w, card_h, (255, 255, 255))
        # 顶部强调色
        add_rectangle(slide, x, y, card_w, Inches(0.08), accent)

        # 顶部色块（用于图标占位）
        icon_area = add_rectangle(slide, x, y + Inches(0.08), card_w, Inches(1.5), primary)

        # 序号圆形
        circle_x = x + card_w/2 - Inches(0.4)
        circle = slide.shapes.add_shape(9, circle_x, y + Inches(0.3), Inches(0.8), Inches(0.8))
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(*accent)
        circle.line.fill.background()

        # 序号文字
        add_text(slide, f"{idx + 1}",
                 circle_x, y + Inches(0.4), Inches(0.8), Inches(0.6),
                 24, (255, 255, 255), True, style["font_title"], PP_ALIGN.CENTER)

        # 内容
        add_text(slide, bullet,
                 x + Inches(0.2), y + Inches(2), card_w - Inches(0.4), Inches(2.5),
                 15, (80, 80, 80), False, style["font_body"], PP_ALIGN.CENTER, 1.5)


# 双栏布局
def add_two_col_slide(prs, data, primary, secondary, accent, style, images):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 顶部色带
    add_rectangle(slide, 0, 0, prs.slide_width, Inches(0.08), accent)
    add_rectangle(slide, 0, 0, prs.slide_width, Inches(1.3), primary)

    add_text(slide, data.get("title", ""),
             Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             30, (255, 255, 255), True, style["font_title"], PP_ALIGN.LEFT)

    # 左右两栏
    bullets = data.get("bullets", [])
    half = len(bullets) // 2 + 1

    left_content = "\n".join([f"• {b}" for b in bullets[:half]])
    right_content = "\n".join([f"• {b}" for b in bullets[half:]])

    # 左栏
    add_rectangle(slide, Inches(0.4), Inches(1.6), Inches(0.06), Inches(4.5), accent)
    add_text(slide, left_content,
             Inches(0.6), Inches(1.6), Inches(5.5), Inches(5),
             18, (51, 51, 51), False, style["font_body"], PP_ALIGN.LEFT, 1.8)

    # 右栏
    add_rectangle(slide, Inches(6.8), Inches(1.6), Inches(0.06), Inches(4.5), secondary)
    add_text(slide, right_content,
             Inches(7), Inches(1.6), Inches(5.5), Inches(5),
             18, (51, 51, 51), False, style["font_body"], PP_ALIGN.LEFT, 1.8)

    # 中间装饰线
    add_rectangle(slide, Inches(6.4), Inches(1.6), Inches(0.02), Inches(4.5), (220, 220, 220))


# 结束页
def add_ending_slide(prs, data, primary, secondary, accent, style):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 渐变背景
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_gradient(bg, secondary, primary, 45)
    bg.line.fill.background()

    # 装饰
    add_decorative_circle(slide, Inches(-2), Inches(-2), Inches(7), accent, 0.1)
    add_decorative_circle(slide, Inches(9), Inches(4), Inches(5), (255, 255, 255), 0.05)

    # 谢谢文字
    add_text(slide, data.get("title", "谢谢"),
             Inches(0), Inches(2.5), prs.slide_width, Inches(1.5),
             60, (255, 255, 255), True, style["font_title"], PP_ALIGN.CENTER)

    if data.get("subtitle"):
        add_text(slide, data.get("subtitle"),
                 Inches(0), Inches(4.2), prs.slide_width, Inches(0.8),
                 24, accent, False, style["font_body"], PP_ALIGN.CENTER)