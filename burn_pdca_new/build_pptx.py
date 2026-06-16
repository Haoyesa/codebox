# -*- coding: utf-8 -*-
"""PPTX builder for the redesigned PDCA report."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

NAVY     = RGBColor(0x0F, 0x2A, 0x5C)
BLUE     = RGBColor(0x1E, 0x4F, 0xB5)
SKY      = RGBColor(0x3B, 0x82, 0xF6)
LIGHT    = RGBColor(0xE8, 0xF0, 0xFE)
ACCENT   = RGBColor(0xC8, 0x1D, 0x25)
GOLD     = RGBColor(0xD4, 0x9B, 0x3B)
INK      = RGBColor(0x1F, 0x29, 0x37)
MUTED    = RGBColor(0x6B, 0x72, 0x80)
RULE     = RGBColor(0xD1, 0xD5, 0xDB)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
PANEL    = RGBColor(0xF6, 0xF8, 0xFC)
SOFT     = RGBColor(0xEE, 0xF2, 0xF9)
ZH_FONT  = "Microsoft YaHei"
EN_FONT  = "Calibri"
SLIDE_W  = Inches(13.333)
SLIDE_H  = Inches(7.5)
TOTAL    = 12

# ---------- Low-level helpers ----------
def add_rect(slide, x, y, w, h, fill, line=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_w is not None:
            shp.line.width = line_w
    shp.shadow.inherit = False
    return shp

def add_round(slide, x, y, w, h, fill, line=None, line_w=None, radius_adj=0.08):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        shp.adjustments[0] = radius_adj
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_w is not None:
            shp.line.width = line_w
    shp.shadow.inherit = False
    return shp

def add_oval(slide, x, y, w, h, fill, line=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_w is not None:
            shp.line.width = line_w
    shp.shadow.inherit = False
    return shp

def add_line(slide, x1, y1, x2, y2, color, weight=Pt(1.0)):
    shp = slide.shapes.add_connector(1, x1, y1, x2, y2)
    shp.line.color.rgb = color
    shp.line.width = weight
    return shp

def add_text(slide, x, y, w, h, text, size=14, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=ZH_FONT,
             line_spacing=1.2):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        rPr = run._r.get_or_add_rPr()
        for tag in ('a:ea', 'a:latin'):
            el = rPr.find(qn(tag))
            if el is None:
                el = etree.SubElement(rPr, qn(tag))
            el.set('typeface', ZH_FONT)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb

def add_rich_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.TOP, line_spacing=1.3):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for (txt, size, bold, color, font) in para:
            run = p.add_run()
            run.text = txt
            run.font.name = font or ZH_FONT
            rPr = run._r.get_or_add_rPr()
            for tag in ('a:ea', 'a:latin'):
                el = rPr.find(qn(tag))
                if el is None:
                    el = etree.SubElement(rPr, qn(tag))
                el.set('typeface', font or ZH_FONT)
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    return tb

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_page_chrome(slide, page_no, total, label='提升理疗防烫伤健康宣教知晓率 · PDCA 持续质量改进'):
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(0.08), NAVY)
    add_line(slide, Inches(0.6), Inches(7.18), Inches(12.733), Inches(7.18), RULE, Pt(0.75))
    add_text(slide, Inches(0.6), Inches(7.22), Inches(8), Inches(0.3),
             label, size=9, color=MUTED)
    add_text(slide, Inches(11.5), Inches(7.22), Inches(1.233), Inches(0.3),
             '%02d / %02d' % (page_no, total), size=9, color=MUTED, align=PP_ALIGN.RIGHT)

def add_section_header(slide, kicker, title):
    add_text(slide, Inches(0.6), Inches(0.5), Inches(6), Inches(0.32),
             kicker, size=11, bold=True, color=BLUE)
    add_text(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.7),
             title, size=26, bold=True, color=NAVY)
    add_rect(slide, Inches(0.6), Inches(1.55), Inches(0.6), Inches(0.06), ACCENT)
    add_rect(slide, Inches(1.25), Inches(1.55), Inches(11.5), Inches(0.02), RULE)


# ---------- Slide 1: Cover ----------
def build_cover(prs):
    s = blank_slide(prs)
    set_bg(s, WHITE)
    add_rect(s, Emu(0), Emu(0), Inches(5.2), SLIDE_H, NAVY)
    add_rect(s, Inches(5.2), Emu(0), Inches(0.12), SLIDE_H, ACCENT)
    add_text(s, Inches(0.7), Inches(1.0), Inches(4.2), Inches(0.4),
             "PDCA · 持续质量改进", size=14, bold=True, color=RGBColor(0xCB, 0xD8, 0xF0))
    add_text(s, Inches(0.7), Inches(1.5), Inches(4.2), Inches(1.3),
             "2026", size=72, bold=True, color=WHITE, font=EN_FONT)
    add_text(s, Inches(0.7), Inches(6.45), Inches(4.2), Inches(0.32),
             "REPORT", size=10, bold=True, color=GOLD)
    add_text(s, Inches(0.7), Inches(6.75), Inches(4.2), Inches(0.32),
             "Quality Improvement Project", size=12, color=RGBColor(0xCB, 0xD8, 0xF0), font=EN_FONT)

    add_text(s, Inches(5.7), Inches(1.6), Inches(7), Inches(0.4),
             "医疗质量改进项目汇报", size=14, bold=True, color=BLUE)
    add_text(s, Inches(5.7), Inches(2.0), Inches(7.2), Inches(2.2),
             ["提升理疗防烫伤", "健康宣教知晓率"], size=38, bold=True, color=NAVY, line_spacing=1.15)
    add_rect(s, Inches(5.7), Inches(4.05), Inches(0.5), Inches(0.06), ACCENT)
    add_text(s, Inches(5.7), Inches(4.2), Inches(7.2), Inches(0.5),
             "PLAN · DO · CHECK · ACT", size=14, bold=True, color=BLUE, font=EN_FONT)
    add_text(s, Inches(5.7), Inches(4.6), Inches(7.2), Inches(0.5),
             "PDCA 循环管理 · 全周期安全升级", size=16, color=INK)
    add_rect(s, Inches(5.7), Inches(5.7), Inches(7.2), Inches(1.0), PANEL)
    add_rect(s, Inches(5.7), Inches(5.7), Inches(0.08), Inches(1.0), BLUE)
    add_rich_text(s, Inches(5.95), Inches(5.78), Inches(7), Inches(0.9), [
        [("汇报科室  ", 11, True,  MUTED, ZH_FONT), ("康复医学科 / 护理部 / 质量管理办公室", 13, False, INK, ZH_FONT)],
        [("汇报日期  ", 11, True,  MUTED, ZH_FONT), ("2026 年 5 月 27 日", 13, False, INK, ZH_FONT)],
    ], line_spacing=1.5)


# ---------- Slide 2: Background ----------
def build_background(prs, page):
    s = blank_slide(prs)
    set_bg(s, WHITE)
    add_page_chrome(s, page, TOTAL)
    add_section_header(s, "01  项目背景", "现状堪忧：理疗相关烫伤事件不容忽视")
    add_text(s, Inches(0.6), Inches(1.85), Inches(12.2), Inches(0.7),
             "理疗（红外线、热敷、蜡疗等）是康复治疗的重要手段，但因操作不规范、设备使用不当或患者自身感觉减退，"
             "导致的烫伤事件屡见不鲜，已成为医疗安全管理中不可忽视的潜在风险。",
             size=13, color=INK, line_spacing=1.55)

    risks = [
        ("风险居高不下", "康复科艾灸及热疗相关烫伤发生率长期处于高位，基层机构与居家场景中安全事故仍时有发生。"),
        ("脆弱群体易感", "糖尿病、脑卒中、脊髓损伤及老年患者因感觉迟钝，对温度感知滞后，是主要受害者。"),
        ("低温烫伤深远", "创面深、愈合慢、易感染，延长住院周期、增加医疗成本，甚至引发医疗纠纷。"),
    ]
    card_w = Inches(3.95)
    gap = Inches(0.2)
    start_x = Inches(0.6)
    top = Inches(2.85)
    for i, (h, body) in enumerate(risks):
        x = start_x + (card_w + gap) * i
        add_round(s, x, top, card_w, Inches(2.2), PANEL, line=RULE, line_w=Pt(0.75))
        add_rect(s, x, top, card_w, Inches(0.45), NAVY)
        add_text(s, x + Inches(0.25), top + Inches(0.07), card_w - Inches(0.5), Inches(0.32),
                 "0%d" % (i+1), size=14, bold=True, color=GOLD, font=EN_FONT)
        add_text(s, x + Inches(0.85), top + Inches(0.07), card_w - Inches(1), Inches(0.32),
                 h, size=14, bold=True, color=WHITE)
        add_text(s, x + Inches(0.25), top + Inches(0.7), card_w - Inches(0.5), Inches(1.4),
                 body, size=12, color=INK, line_spacing=1.5)

    case_top = Inches(5.25)
    add_rect(s, Inches(0.6), case_top, Inches(0.08), Inches(1.55), ACCENT)
    add_text(s, Inches(0.85), case_top, Inches(12), Inches(0.32),
             "真实警示案例", size=12, bold=True, color=ACCENT)
    add_rich_text(s, Inches(0.85), case_top + Inches(0.35), Inches(12.0), Inches(1.2), [
        [("案例一  ", 12, True, NAVY, ZH_FONT),
         ("62 岁老年患者自行用吹风机高温直吹右踝，造成 Ⅲ 度低温烫伤，住院 32 天。", 12, False, INK, ZH_FONT)],
        [("案例二  ", 12, True, NAVY, ZH_FONT),
         ("糖尿病患者无指导下使用理疗仪处理富贵包，颈背出现 10×10 cm 深度溃疡，愈合极困难。", 12, False, INK, ZH_FONT)],
    ], line_spacing=1.5)


# ---------- Slide 3: Problem ----------
def build_problem(prs, page):
    s = blank_slide(prs)
    set_bg(s, WHITE)
    add_page_chrome(s, page, TOTAL)
    add_section_header(s, "02  问题分析", "核心痛点：宣教不足与认知缺失")

    add_text(s, Inches(0.6), Inches(1.85), Inches(5.9), Inches(0.4),
             "对标高标准：医疗规范的刚性要求", size=14, bold=True, color=BLUE)
    items_l = [
        ("评审标准红线", "保障患者安全，从制度源头减少因宣教缺位导致的不良临床事件。"),
        ("安全目标导向", "强调医患有效沟通，确保患方清晰理解康复方案的风险与注意事项。"),
        ("专业指南遵循", "落实《康复治疗安全风险防控指南》，必须执行充分的安全教育流程。"),
    ]
    top = Inches(2.35)
    for i, (h, body) in enumerate(items_l):
        y = top + Inches(1.2) * i
        add_rect(s, Inches(0.6), y, Inches(0.05), Inches(1.0), BLUE)
        add_text(s, Inches(0.8), y, Inches(5.8), Inches(0.35),
                 h, size=13, bold=True, color=NAVY)
        add_text(s, Inches(0.8), y + Inches(0.38), Inches(5.8), Inches(0.7),
                 body, size=11, color=INK, line_spacing=1.4)

    add_text(s, Inches(7.0), Inches(1.85), Inches(5.8), Inches(0.4),
             "现实挑战：执行层面的落地断层", size=14, bold=True, color=ACCENT)
    items_r = [
        ("宣教形式僵化", "依赖口头告知，缺乏标准化、可视化的图文或视频材料。"),
        ("介入时机错位", "常在治疗繁忙时段突击宣教，信息接收效果差。"),
        ("患者认知盲区", "对低温烫伤机制理解不足，存在侥幸心理。"),
        ("效果闭环断裂", "缺乏量化评估工具，无法确认是否真正掌握。"),
    ]
    top = Inches(2.35)
    for i, (h, body) in enumerate(items_r):
        y = top + Inches(0.92) * i
        add_round(s, Inches(7.0), y, Inches(5.9), Inches(0.78), PANEL, line=RULE, line_w=Pt(0.5), radius_adj=0.2)
        add_text(s, Inches(7.2), y + Inches(0.08), Inches(1.6), Inches(0.32),
             "×  " + h, size=12, bold=True, color=ACCENT)
        add_text(s, Inches(8.7), y + Inches(0.1), Inches(4.0), Inches(0.6),
                 body, size=11, color=INK, line_spacing=1.35)

    add_round(s, Inches(0.6), Inches(6.2), Inches(12.3), Inches(0.85), LIGHT, line=None, radius_adj=0.2)
    add_text(s, Inches(0.95), Inches(6.32), Inches(11.8), Inches(0.6),
             "↳  行业规范已将安全教育划定清晰底线，必须将 “被动告知” 转变为 “主动保障”。",
             size=12, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)


# ---------- Slide 4: Root cause (4M) ----------
def build_rootcause(prs, page):
    s = blank_slide(prs)
    set_bg(s, WHITE)
    add_page_chrome(s, page, TOTAL)
    add_section_header(s, "03  根因分析", "追根溯源：理疗烫伤为何发生？")

    grid = [
        ("M · 人员", "Man", [
            "宣教意识不强、培训不规范、评估流于形式",
            "患方健康素养差异，存侥幸心理",
        ]),
        ("M · 方法", "Method", [
            "形式单一，仅靠口头告知",
            "缺图文/视频，缺量化评估问卷",
        ]),
        ("M · 环境", "Environment", [
            "治疗区繁忙，缺独立安静宣教空间",
            "嘈杂影响沟通效率与信息完整度",
        ]),
        ("M · 物料", "Material", [
            "缺统一规范的防烫伤宣教手册/视频",
            "专业性过强，缺乏生活化案例",
        ]),
    ]
    col_w = Inches(6.05)
    row_h = Inches(2.05)
    x0, y0 = Inches(0.6), Inches(1.95)
    for i, (label, en, items) in enumerate(grid):
        r, c = i // 2, i % 2
        x = x0 + (col_w + Inches(0.2)) * c
        y = y0 + (row_h + Inches(0.18)) * r
        add_round(s, x, y, col_w, row_h, WHITE, line=RULE, line_w=Pt(0.75), radius_adj=0.05)
        add_rect(s, x, y, Inches(1.4), row_h, NAVY)
        add_text(s, x, y + Inches(0.4), Inches(1.4), Inches(0.6),
                 label, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(0.95), Inches(1.4), Inches(0.4),
                 en, size=12, color=GOLD, font=EN_FONT, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(1.6), y + Inches(0.25), col_w - Inches(1.8), Inches(0.6),
                 "· " + items[0], size=12, color=INK, line_spacing=1.4)
        add_text(s, x + Inches(1.6), y + Inches(0.95), col_w - Inches(1.8), Inches(0.6),
                 "· " + items[1], size=12, color=INK, line_spacing=1.4)

    insight_top = Inches(6.25)
    add_rect(s, Inches(0.6), insight_top, Inches(12.3), Inches(0.8), NAVY)
    add_text(s, Inches(0.95), insight_top + Inches(0.12), Inches(11.8), Inches(0.32),
             "核心症结洞察", size=11, bold=True, color=GOLD)
    add_text(s, Inches(0.95), insight_top + Inches(0.42), Inches(11.8), Inches(0.32),
             "宣教形式单一 + 效果评估缺失 = 风险闭环断裂；需从标准化工具、流程优化、环境改造三方面同步发力。",
             size=12, color=WHITE)


# ---------- Slide 5: SMART goal ----------
def build_goal(prs, page):
    s = blank_slide(prs)
    set_bg(s, WHITE)
    add_page_chrome(s, page, TOTAL)
    add_section_header(s, "04  目标设定", "明确目标：基于 SMART 原则的安全升级")

    add_text(s, Inches(0.6), Inches(1.85), Inches(12.2), Inches(0.85),
             "通过 PDCA 循环管理，系统提升理疗患者对防烫伤核心知识的全面知晓率，建立标准化安全防护流程，"
             "显著降低理疗相关烫伤（特别是低温烫伤）的发生率，实现 “零非预期烫伤” 的安全愿景。",
             size=12, color=INK, line_spacing=1.5)

    smarts = [
        ("S", "Specific", "聚焦核心", ["清晰掌握安全注意事项", "重点锁定低温烫伤", "制定专项预防与处置标准"]),
        ("M", "Measurable", "可量化",   ["知晓率 65% → 95%+", "烫伤不良事件下降 50%", "问卷合格率量化考核"]),
        ("A", "Achievable", "可达成",   ["依托现有医疗资源", "图文 SOP + 多媒体材料", "全员专项操作培训"]),
        ("R", "Relevant",  "相关价值", ["以患者安全为中心", "符合等级医院评审规范", "提升服务质量与运营安全"]),
        ("T", "Time-bound", "有时限",   ["3 个月完整 PDCA 循环", "明确阶段里程碑", "按时保质达成目标"]),
    ]
    col_w = Inches(2.42)
    gap   = Inches(0.1)
    x0    = Inches(0.6)
    y0    = Inches(2.85)
    for i, (letter, en, cn, items) in enumerate(smarts):
        x = x0 + (col_w + gap) * i
        add_rect(s, x, y0, col_w, Inches(0.75), NAVY)
        add_text(s, x, y0 + Inches(0.05), col_w, Inches(0.45),
                 letter, size=26, bold=True, color=GOLD, font=EN_FONT, align=PP_ALIGN.CENTER)
        add_text(s, x, y0 + Inches(0.5), col_w, Inches(0.25),
                 en, size=10, color=WHITE, font=EN_FONT, align=PP_ALIGN.CENTER)
        add_rect(s, x, y0 + Inches(0.75), col_w, Inches(2.3), PANEL, line=RULE, line_w=Pt(0.5))
        add_text(s, x + Inches(0.2), y0 + Inches(0.85), col_w - Inches(0.4), Inches(0.3),
                 cn, size=12, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        for j, it in enumerate(items):
            add_text(s, x + Inches(0.2), y0 + Inches(1.2) + Inches(0.6) * j,
                     col_w - Inches(0.4), Inches(0.6),
                     "· " + it, size=10.5, color=INK, line_spacing=1.25)

    kpi_y = Inches(6.2)
    add_rect(s, Inches(0.6), kpi_y, Inches(12.3), Inches(0.85), LIGHT, line=None)
    kpis = [("知晓率目标", "95%+"), ("不良事件下降", "50%"), ("执行周期", "3 个月"), ("培训覆盖率", "100%")]
    seg_w = Inches(12.3 / 4)
    for i, (label, val) in enumerate(kpis):
        x = Inches(0.6) + seg_w * i
        if i > 0:
            add_line(s, x, kpi_y + Inches(0.1), x, kpi_y + Inches(0.75), RULE, Pt(0.5))
        add_text(s, x, kpi_y + Inches(0.12), seg_w, Inches(0.3),
                 label, size=10, color=MUTED, align=PP_ALIGN.CENTER)
        add_text(s, x, kpi_y + Inches(0.4), seg_w, Inches(0.45),
                 val, size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ---------- Slide 6: Strategy ----------
def build_strategy(prs, page):
    s = blank_slide(prs)
    set_bg(s, WHITE)
    add_page_chrome(s, page, TOTAL)
    add_section_header(s, "05  改进策略", "群策群力：四大改进方向协同发力")

    strategies = [
        ("01", "宣教内容与形式", "标准化工具包 + 多元化互动", [
            "图文手册 / 科普短视频 / 情景卡片",
            "口头 + 视频 + 实操 + 提问 四维模式",
            "专业知识 → 患者易懂、易记、易用",
        ]),
        ("02", "宣教流程与时机", "SOP 化操作 + 黄金窗口", [
            "《理疗防烫伤健康宣教 SOP》",
            "治疗前 15-20 分钟黄金干预",
            "独立区域，一对一专项指导",
        ]),
        ("03", "医护人员培训", "全员赋能 + 绩效闭环", [
            "新版 SOP + 跨年龄沟通技巧",
            "风险快速评估方法轮训",
            "宣教质量纳入绩效考核",
        ]),
        ("04", "效果评估机制", "即时验证 + 持续改进", [
            "《防烫伤知识知晓率问卷》",
            "宣教后与治疗结束前测评",
            "执行 - 评估 - 优化 PDCA 闭环",
        ]),
    ]
    col_w = Inches(6.05)
    row_h = Inches(2.3)
    x0, y0 = Inches(0.6), Inches(1.85)
    for i, (no, h, sub, items) in enumerate(strategies):
        r, c = i // 2, i % 2
        x = x0 + (col_w + Inches(0.2)) * c
        y = y0 + (row_h + Inches(0.2)) * r
        add_round(s, x, y, col_w, row_h, WHITE, line=RULE, line_w=Pt(0.75), radius_adj=0.05)
        add_rect(s, x, y, Inches(1.0), row_h, NAVY)
        add_text(s, x, y + Inches(0.7), Inches(1.0), Inches(0.7),
                 no, size=32, bold=True, color=GOLD, font=EN_FONT, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(1.45), Inches(1.0), Inches(0.3),
                 "STRATEGY", size=8, color=RGBColor(0xCB, 0xD8, 0xF0), font=EN_FONT, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(1.2), y + Inches(0.2), col_w - Inches(1.4), Inches(0.4),
                 h, size=15, bold=True, color=NAVY)
        add_text(s, x + Inches(1.2), y + Inches(0.6), col_w - Inches(1.4), Inches(0.32),
                 sub, size=11, color=BLUE, bold=True)
        for j, it in enumerate(items):
            add_text(s, x + Inches(1.2), y + Inches(1.0) + Inches(0.4) * j,
                     col_w - Inches(1.4), Inches(0.4),
                     "· " + it, size=10.5, color=INK, line_spacing=1.3)


# ---------- Slide 7: Roadmap ----------
def build_roadmap(prs, page):
    s = blank_slide(prs)
    set_bg(s, WHITE)
    add_page_chrome(s, page, TOTAL)
    add_section_header(s, "06  实施路径", "明确分工与时间表：把措施落到实处")

    add_text(s, Inches(0.6), Inches(1.85), Inches(6.0), Inches(0.4),
             "组织保障 · 关键角色", size=14, bold=True, color=BLUE)
    roles = [
        ("项目负责人", "统筹规划、跨部门协调、对落地效果负总责"),
        ("康复科护士长", "主导制定 SOP，组织培训，监督执行"),
        ("康复治疗师",  "参与方案开发，融入安全要点提醒"),
        ("护理骨干",    "宣教包一线实施，问卷发放回收"),
        ("质控科干事",  "数据收集、统计分析、效果评估"),
        ("信息科 / 宣传科", "视频拍摄、图文设计、可视化呈现"),
    ]
    top = Inches(2.35)
    for i, (h, body) in enumerate(roles):
        y = top + Inches(0.62) * i
        add_rect(s, Inches(0.6), y + Inches(0.05), Inches(0.18), Inches(0.18), BLUE)
        add_text(s, Inches(0.9), y, Inches(1.7), Inches(0.32),
                 h, size=11, bold=True, color=NAVY)
        add_text(s, Inches(2.55), y, Inches(4.1), Inches(0.32),
                 body, size=10.5, color=INK)

    add_text(s, Inches(7.0), Inches(1.85), Inches(5.8), Inches(0.4),
             "12 周推进计划", size=14, bold=True, color=BLUE)

    # timeline vertical column
    phases = [
        ("W1",     "筹备期",  "方案细化、宣教视频与图文物料制作"),
        ("W2-3",   "赋能期",  "全员培训、发布执行 SOP、统一标准"),
        ("W4-10",  "执行期",  "新流程全面落地、一线规范执行"),
        ("W11-12", "结项期",  "数据复盘、效果评估、输出报告"),
    ]
    y0 = Inches(2.4)
    row_h = Inches(1.0)
    for i, (wk, name, desc) in enumerate(phases):
        y = y0 + row_h * i
        add_oval(s, Inches(7.05), y + Inches(0.05), Inches(0.4), Inches(0.4), ACCENT, line=WHITE, line_w=Pt(1.5))
        add_text(s, Inches(7.05), y + Inches(0.05), Inches(0.4), Inches(0.4),
                 str(i+1), size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(7.6), y, Inches(1.4), Inches(0.3),
                 wk, size=12, bold=True, color=BLUE, font=EN_FONT)
        add_text(s, Inches(7.6), y + Inches(0.3), Inches(1.4), Inches(0.3),
                 name, size=11, bold=True, color=NAVY)
        add_text(s, Inches(7.6), y + Inches(0.6), Inches(5.5), Inches(0.4),
                 desc, size=10, color=INK, line_spacing=1.3)
        if i < len(phases) - 1:
            add_line(s, Inches(7.25), y + Inches(0.45), Inches(7.25), y + row_h, RULE, Pt(1.0))

    add_round(s, Inches(0.6), Inches(6.3), Inches(12.3), Inches(0.78), LIGHT, line=None, radius_adj=0.25)
    add_text(s, Inches(0.95), Inches(6.42), Inches(11.8), Inches(0.55),
             "↳  全员培训覆盖率 100% / 现场督导整改率 100%，把 “防烫伤” 融入每一次理疗服务。",
             size=12, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)


# ---------- Slide 8: Challenges ----------
def build_challenges(prs, page):
    s = blank_slide(prs)
    set_bg(s, WHITE)
    add_page_chrome(s, page, TOTAL)
    add_section_header(s, "07  问题应对", "攻坚克难：把阻力变成改进契机")

    challenges = [
        ("人员协作难题", "项目初期工作量增加，部分护士因节奏改变产生抵触情绪。",
         "管理优化 + 激励并行", "加强沟通疏导，优化排班；宣教质量纳入月度绩效考核，正向激励。"),
        ("老年沟通障碍", "目标人群中大量老年患者视力不佳，常规手册字体偏小，理解度差。",
         "适老化 + 人性化", "紧急改版大字版手册；增加一对一讲解环节，放慢语速、耐心解答。"),
        ("设备设施限制", "宣教视频需配合治疗同步观看，但部分治疗室缺乏固定播放设备。",
         "移动化硬件升级", "协调信息科为每个治疗室配备可移动触控平板，适配流动服务场景。"),
    ]
    y0 = Inches(1.85)
    for i, (h, prob, sub, sol) in enumerate(challenges):
        y = y0 + Inches(1.55) * i
        add_round(s, Inches(0.6), y, Inches(12.3), Inches(1.4), PANEL, line=RULE, line_w=Pt(0.5), radius_adj=0.05)
        add_rect(s, Inches(0.6), y, Inches(0.15), Inches(1.4), ACCENT)
        add_text(s, Inches(0.95), y + Inches(0.1), Inches(2.6), Inches(0.32),
                 "挑战 %d" % (i+1), size=10, bold=True, color=ACCENT)
        add_text(s, Inches(0.95), y + Inches(0.4), Inches(2.6), Inches(0.4),
                 h, size=14, bold=True, color=NAVY)
        add_text(s, Inches(0.95), y + Inches(0.85), Inches(2.6), Inches(0.55),
                 prob, size=10.5, color=INK, line_spacing=1.3)
        add_text(s, Inches(3.7), y + Inches(0.5), Inches(0.4), Inches(0.4),
                 "→", size=22, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        add_text(s, Inches(4.2), y + Inches(0.1), Inches(8.6), Inches(0.32),
                 "破局对策  ·  " + sub, size=11, bold=True, color=BLUE)
        add_text(s, Inches(4.2), y + Inches(0.45), Inches(8.6), Inches(0.95),
                 sol, size=12, color=INK, line_spacing=1.5)


# ---------- Slide 9: Results ----------
def build_results(prs, page):
    s = blank_slide(prs)
    set_bg(s, WHITE)
    add_page_chrome(s, page, TOTAL)
    add_section_header(s, "08  改进成效", "成果显著：核心指标全面改善")

    kpi_y = Inches(1.95)
    add_round(s, Inches(0.6), kpi_y, Inches(6.05), Inches(2.6), WHITE, line=BLUE, line_w=Pt(1.5), radius_adj=0.04)
    add_rect(s, Inches(0.6), kpi_y, Inches(6.05), Inches(0.5), NAVY)
    add_text(s, Inches(0.85), kpi_y + Inches(0.1), Inches(5.5), Inches(0.32),
             "患者安全知识知晓率", size=12, bold=True, color=WHITE)
    add_text(s, Inches(0.85), kpi_y + Inches(0.75), Inches(2.0), Inches(0.4),
             "改善前", size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.85), kpi_y + Inches(1.1), Inches(2.0), Inches(0.9),
             "65%", size=44, bold=True, color=MUTED, font=EN_FONT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(3.6), kpi_y + Inches(0.75), Inches(2.0), Inches(0.4),
             "改善后", size=11, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(3.6), kpi_y + Inches(1.1), Inches(2.0), Inches(0.9),
             "96.2%", size=44, bold=True, color=ACCENT, font=EN_FONT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(5.5), kpi_y + Inches(1.05), Inches(1.05), Inches(0.9),
             "+31.2pp", size=16, bold=True, color=GOLD, font=EN_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.85), kpi_y + Inches(2.05), Inches(5.5), Inches(0.5),
             "合格率大幅提升，远超 95% 预设目标，标志患者理解已达理想水平。",
             size=11, color=INK, line_spacing=1.4)

    add_round(s, Inches(6.85), kpi_y, Inches(6.05), Inches(2.6), WHITE, line=ACCENT, line_w=Pt(1.5), radius_adj=0.04)
    add_rect(s, Inches(6.85), kpi_y, Inches(6.05), Inches(0.5), ACCENT)
    add_text(s, Inches(7.1), kpi_y + Inches(0.1), Inches(5.5), Inches(0.32),
             "理疗相关烫伤不良事件", size=12, bold=True, color=WHITE)
    add_text(s, Inches(7.1), kpi_y + Inches(0.75), Inches(2.5), Inches(0.4),
             "同期历史发生数", size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, Inches(7.1), kpi_y + Inches(1.1), Inches(2.5), Inches(0.9),
             "X 例", size=44, bold=True, color=MUTED, font=EN_FONT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(9.9), kpi_y + Inches(0.75), Inches(2.5), Inches(0.4),
             "本期实际发生数", size=11, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(9.9), kpi_y + Inches(1.1), Inches(2.5), Inches(0.9),
             "0 例", size=44, bold=True, color=ACCENT, font=EN_FONT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(7.1), kpi_y + Inches(2.05), Inches(5.5), Inches(0.5),
             "标准化流程 + 主动干预，消除烫伤隐患，实现阶段性零事故管理目标。",
             size=11, color=INK, line_spacing=1.4)

    bar_y = Inches(4.85)
    add_text(s, Inches(0.6), bar_y, Inches(12.0), Inches(0.35),
             "关键发现：安全距离 / 禁止自行调节温度 等知识点掌握度提升最显著",
             size=12, bold=True, color=NAVY)
    add_text(s, Inches(0.6), bar_y + Inches(0.45), Inches(2), Inches(0.32),
             "改善前", size=11, color=MUTED)
    bar_x = Inches(2.0)
    bar_max = Inches(8.0)
    add_rect(s, bar_x, bar_y + Inches(0.45), bar_max, Inches(0.4), RULE)
    add_rect(s, bar_x, bar_y + Inches(0.45), Emu(int(bar_max * 0.65)), Inches(0.4), MUTED)
    add_text(s, bar_x + Emu(int(bar_max * 0.65)) + Inches(0.1), bar_y + Inches(0.45), Inches(1.2), Inches(0.4),
             "65%", size=11, bold=True, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.6), bar_y + Inches(1.0), Inches(2), Inches(0.32),
             "改善后", size=11, color=ACCENT, bold=True)
    add_rect(s, bar_x, bar_y + Inches(1.0), bar_max, Inches(0.4), LIGHT)
    add_rect(s, bar_x, bar_y + Inches(1.0), Emu(int(bar_max * 0.962)), Inches(0.4), ACCENT)
    add_text(s, bar_x + Emu(int(bar_max * 0.962)) + Inches(0.1), bar_y + Inches(1.0), Inches(1.2), Inches(0.4),
             "96.2%", size=11, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)


# ---------- Slide 10: Bonus ----------
def build_bonus(prs, page):
    s = blank_slide(prs)
    set_bg(s, WHITE)
    add_page_chrome(s, page, TOTAL)
    add_section_header(s, "09  意外收获", "患者满意度与护理效率双提升")

    add_round(s, Inches(0.6), Inches(1.85), Inches(5.4), Inches(4.85), NAVY, line=None, radius_adj=0.04)
    add_text(s, Inches(0.85), Inches(2.05), Inches(5.0), Inches(0.4),
             "患者满意度", size=12, bold=True, color=GOLD)
    add_text(s, Inches(0.85), Inches(2.5), Inches(5.0), Inches(1.2),
             "92%  →  98%", size=42, bold=True, color=WHITE, font=EN_FONT)
    add_text(s, Inches(0.85), Inches(3.7), Inches(5.0), Inches(0.4),
             "服务认可度从合格迈向卓越", size=13, color=RGBColor(0xCB, 0xD8, 0xF0))
    add_text(s, Inches(0.85), Inches(4.1), Inches(5.0), Inches(0.4),
             "患者对医疗安全的感知度与信任度同步飞跃", size=11, color=RGBColor(0xCB, 0xD8, 0xF0))
    add_rect(s, Inches(0.85), Inches(4.65), Inches(4.95), Inches(0.02), GOLD)
    add_text(s, Inches(0.85), Inches(4.75), Inches(4.95), Inches(0.7),
             "“宣教视频一看就懂，感觉医院把我们的安全真正放在了心上，比光说管用多了！”",
             size=10, color=RGBColor(0xE6, 0xEE, 0xFB), line_spacing=1.4)
    add_text(s, Inches(0.85), Inches(5.55), Inches(4.95), Inches(0.4),
             "— 患者反馈", size=9, color=GOLD, align=PP_ALIGN.RIGHT)
    add_rect(s, Inches(0.85), Inches(6.1), Inches(4.95), Inches(0.02), GOLD)
    add_text(s, Inches(0.85), Inches(6.2), Inches(4.95), Inches(0.4),
             "护理效率 / 安全理念 / 协同闭环同步改善", size=10, color=WHITE, bold=True)

    benefits = [
        ("流程标准化落地", "SOP 消除沟通歧义，压缩重复答疑耗时，让护理节奏更高效。",
         "“标准化的工具让我们的工作更从容，专业形象和自信心都得到了很大的提升。”"),
        ("安全理念深化", "“安全第一” 融入日常诊疗细节，团队风险防范意识显著增强。",
         "“现在大家更愿意主动发现隐患、上报问题，因为我们知道这是在守护患者和自己。”"),
        ("协同闭环形成", "医护患三方共同参与安全防线，构建可持续的质量改进生态。",
         "“以患者安全为中心” 成为全员共识，注入高质量发展内在动力。"),
    ]
    y0 = Inches(1.85)
    for i, (h, body, quote) in enumerate(benefits):
        y = y0 + Inches(1.6) * i
        add_round(s, Inches(6.2), y, Inches(6.7), Inches(1.45), PANEL, line=RULE, line_w=Pt(0.5), radius_adj=0.05)
        add_rect(s, Inches(6.2), y, Inches(0.12), Inches(1.45), BLUE)
        add_text(s, Inches(6.45), y + Inches(0.12), Inches(6.3), Inches(0.32),
                 h, size=13, bold=True, color=NAVY)
        add_text(s, Inches(6.45), y + Inches(0.5), Inches(6.3), Inches(0.4),
                 body, size=10.5, color=INK, line_spacing=1.4)
        add_text(s, Inches(6.45), y + Inches(0.9), Inches(6.3), Inches(0.5),
                 quote, size=10, color=BLUE, line_spacing=1.4)


# ---------- Slide 11: Sustain ----------
def build_sustain(prs, page):
    s = blank_slide(prs)
    set_bg(s, WHITE)
    add_page_chrome(s, page, TOTAL)
    add_section_header(s, "10  持续改进", "固化成果：让安全成为日常习惯")

    items = [
        ("01", "形成标准化文件", "将 SOP 纳入科室常规制度\n统一问卷为标准评估工具", NAVY),
        ("02", "纳入常规培训",   "新护士岗前培训核心模块\n在职护士年度继续教育", BLUE),
        ("03", "建立长效监督",   "护士长每周现场抽查\n质控科季度专项督查", ACCENT),
        ("04", "推广成功经验",   "形成典型案例全院分享\n辐射更多临床科室", GOLD),
    ]
    y0 = Inches(1.9)
    item_h = Inches(2.4)
    item_w = Inches(6.05)
    for i, (no, h, body, color) in enumerate(items):
        r, c = i // 2, i % 2
        x = Inches(0.6) + (item_w + Inches(0.2)) * c
        y = y0 + (item_h + Inches(0.2)) * r
        add_round(s, x, y, item_w, item_h, WHITE, line=RULE, line_w=Pt(0.75), radius_adj=0.04)
        add_text(s, x + Inches(0.3), y + Inches(0.2), Inches(2), Inches(0.9),
                 no, size=42, bold=True, color=color, font=EN_FONT)
        add_text(s, x + Inches(2.3), y + Inches(0.35), item_w - Inches(2.6), Inches(0.4),
                 h, size=16, bold=True, color=NAVY)
        add_rect(s, x + Inches(2.3), y + Inches(0.85), Inches(1.0), Inches(0.04), color)
        add_text(s, x + Inches(2.3), y + Inches(1.05), item_w - Inches(2.6), Inches(1.3),
                 body, size=11, color=INK, line_spacing=1.5)


# ---------- Slide 12: Reflection + Thanks (combined closer) ----------
def build_reflection(prs, page):
    s = blank_slide(prs)
    set_bg(s, WHITE)
    add_page_chrome(s, page, TOTAL)
    add_section_header(s, "11  反思展望", "反思不足，规划下一阶段")

    add_text(s, Inches(0.6), Inches(1.85), Inches(6.0), Inches(0.4),
             "不足与反思", size=14, bold=True, color=ACCENT)
    refl = [
        ("宣教覆盖有限", "门诊理疗患者的健康宣教环节仍存在空白，需拓宽触达面。"),
        ("形式亟待创新", "口头 + 纸质手册为主，可探索 VR / 互动式线上平台。"),
        ("数据深度不足", "样本时间跨度和规模有限，需长周期持续监测验证。"),
    ]
    top = Inches(2.35)
    for i, (h, body) in enumerate(refl):
        y = top + Inches(1.3) * i
        add_round(s, Inches(0.6), y, Inches(6.0), Inches(1.15), PANEL, line=RULE, line_w=Pt(0.5), radius_adj=0.06)
        add_rect(s, Inches(0.6), y, Inches(0.12), Inches(1.15), ACCENT)
        add_text(s, Inches(0.9), y + Inches(0.12), Inches(5.6), Inches(0.4),
                 h, size=13, bold=True, color=NAVY)
        add_text(s, Inches(0.9), y + Inches(0.5), Inches(5.6), Inches(0.6),
                 body, size=11, color=INK, line_spacing=1.45)

    add_text(s, Inches(7.0), Inches(1.85), Inches(5.8), Inches(0.4),
             "下一阶段核心目标", size=14, bold=True, color=BLUE)
    add_text(s, Inches(7.0), Inches(2.25), Inches(5.8), Inches(0.5),
             "将成熟模式推广至门诊理疗，搭建线上宣教平台，"
             "实现健康管理的长效化与数字化。", size=11, color=INK, line_spacing=1.45)
    plans = [
        ("流程轻量化", "为门诊快节奏场景定制极简宣教流程与配套材料。"),
        ("数字化赋能", "联合信息科开发 APP 专属宣教模块，嵌入智能指引与风险预警。"),
        ("动态追踪",   "建立长期数据档案，持续收集反馈并迭代方案。"),
    ]
    top = Inches(3.3)
    for i, (h, body) in enumerate(plans):
        y = top + Inches(1.0) * i
        add_oval(s, Inches(7.0), y + Inches(0.05), Inches(0.35), Inches(0.35), BLUE)
        add_text(s, Inches(7.0), y + Inches(0.05), Inches(0.35), Inches(0.35),
                 str(i+1), size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(7.5), y, Inches(5.3), Inches(0.4),
                 h, size=13, bold=True, color=NAVY)
        add_text(s, Inches(7.5), y + Inches(0.4), Inches(5.3), Inches(0.5),
                 body, size=11, color=INK, line_spacing=1.4)

    # bottom: thanks band
    add_round(s, Inches(0.6), Inches(6.3), Inches(12.3), Inches(0.78), NAVY, line=None, radius_adj=0.2)
    add_text(s, Inches(0.95), Inches(6.42), Inches(11.8), Inches(0.55),
             "从 “经验” 走向 “标准”，以患者安全为中心的永续改进  ·  Thank You · 期待指导",
             size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


# ---------- Main ----------
def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    build_cover(prs)
    build_background(prs, 2)
    build_problem(prs, 3)
    build_rootcause(prs, 4)
    build_goal(prs, 5)
    build_strategy(prs, 6)
    build_roadmap(prs, 7)
    build_challenges(prs, 8)
    build_results(prs, 9)
    build_bonus(prs, 10)
    build_sustain(prs, 11)
    build_reflection(prs, 12)

    out = r"D:\project\burn_pdca_new\提升理疗防烫伤健康宣教知晓率PDCA汇报_新版.pptx"
    prs.save(out)
    print("Saved:", out)
    print("Slide count:", len(prs.slides.__iter__.__self__._sldIdLst))


if __name__ == "__main__":
    main()



